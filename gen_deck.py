# -*- coding: utf-8 -*-
"""產出：① 深色 PPTX（20 頁，含圖 + 每頁備註）② 上台講稿 Markdown。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = r"D:\07_Claude\01_Master_Report\obesity_project_3"
IMG = os.path.join(BASE, "demo_site", "assets", "img")
PLOT = os.path.join(BASE, "demo_site", "assets", "plots")
OUT_PPTX = os.path.join(BASE, "outputs", "簡報_第五組.pptx")
OUT_MD = os.path.join(BASE, "outputs", "上台講稿_第五組.md")

FONT = "Microsoft JhengHei"
BG = RGBColor(0x0B, 0x11, 0x20)
CARD = RGBColor(0x16, 0x21, 0x3E)
INK = RGBColor(0xE8, 0xEE, 0xFC)
MUTE = RGBColor(0x9A, 0xA7, 0xC8)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
GREEN = RGBColor(0x34, 0xD3, 0x99)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def _set_font(run, size, color, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    # 東亞字型
    rPr = run._r.get_or_add_rPr()
    import copy
    from pptx.oxml.ns import qn
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_bg(slide):
    r = slide.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def tb(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return box, tf


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT, space=6, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    p.line_spacing = 1.15
    runs = text if isinstance(text, list) else [(text, color, bold)]
    for seg in runs:
        r = p.add_run()
        r.text = seg[0]
        _set_font(r, size, seg[1], seg[2] if len(seg) > 2 else bold)
    return p


def header(slide, title_runs, sub=None):
    box, tf = tb(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.2))
    para(tf, title_runs, 30, INK, bold=True, first=True)
    if sub:
        para(tf, sub, 15, MUTE, space=0)
    # accent line
    ln = slide.shapes.add_shape(1, Inches(0.62), Inches(1.5), Inches(2.2), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background(); ln.shadow.inherit = False


def card(slide, l, t, w, h, head, body, accent=BLUE):
    box = slide.shapes.add_shape(1, l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = CARD
    box.line.color.rgb = RGBColor(0x2A, 0x3A, 0x63); box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.12)
    para(tf, head, 16, accent, bold=True, first=True, space=5)
    for b in (body if isinstance(body, list) else [body]):
        para(tf, b, 12.5, INK if False else RGBColor(0xC3,0xCD,0xEE), space=3)
    return box


def add_img_fit(slide, path, l, t, w, h):
    if not os.path.exists(path):
        return None
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = (w / h)
    if ar > box_ar:
        nw = w; nh = int(w / ar)
    else:
        nh = h; nw = int(h * ar)
    nl = l + (w - nw) // 2
    nt = t + (h - nh) // 2
    return slide.shapes.add_picture(path, nl, nt, nw, nh)


def notes(slide, terms, script):
    tf = slide.notes_slide.notes_text_frame
    tf.text = "【名詞解釋】"
    for term, exp in terms:
        p = tf.add_paragraph(); p.text = f"· {term}：{exp}"
    p = tf.add_paragraph(); p.text = ""
    p = tf.add_paragraph(); p.text = "【講稿】"
    for line in script:
        p = tf.add_paragraph(); p.text = line


MD = ["# 上台講稿 · 機器學習期末報告 · 第五組",
      "> 結合生活指標與行為趨勢之肥胖風險預測模型 · 對象：碩班同學",
      "> 每頁含【名詞解釋】(可先簡述) 與【講稿】(精準口白)。\n", "---\n"]


def md_slide(n, title, terms, script):
    MD.append(f"## 第 {n} 頁 · {title}\n")
    if terms:
        MD.append("**名詞解釋**")
        for t, e in terms:
            MD.append(f"- **{t}**：{e}")
        MD.append("")
    MD.append("**講稿**")
    for line in script:
        MD.append(f"> {line}")
    MD.append("\n---\n")


def build(n, title, terms, script, builder):
    s = prs.slides.add_slide(blank)
    add_bg(s)
    builder(s)
    notes(s, terms, script)
    md_slide(n, title, terms, script)
    return s


# ============ 內容 ============
def s1(s):
    box, tf = tb(s, Inches(0.7), Inches(2.1), Inches(7.6), Inches(3.6))
    para(tf, "機器學習期末報告 · 第五組", 16, MUTE, first=True, space=12)
    para(tf, [("結合生活指標與行為趨勢之", INK, True)], 36, INK, bold=True, space=2)
    para(tf, [("肥胖風險預測模型", BLUE, True)], 36, BLUE, bold=True, space=14)
    para(tf, "14773008 賴百擎 · 11773037 陳宗傑", 15, MUTE, space=14)
    para(tf, "不靠醫療紀錄、不靠 BMI，只用日常生活行為，能否準確預測肥胖風險？最少要幾個變數？", 15, AMBER, space=0)
    add_img_fit(s, os.path.join(IMG, "cover_hero.png"), Inches(8.5), Inches(1.6), Inches(4.3), Inches(4.3))

build(1, "封面",
      [("BMI", "身體質量指數＝體重(kg)÷身高(m)²，臨床用來分類過重/肥胖的指標"),
       ("特徵 / 變數", "模型用來預測的輸入欄位，例如睡眠時數、是否運動")],
      ["大家好，我們是第五組。我們的題目是『結合生活指標與行為趨勢之肥胖風險預測模型』。",
       "一句話講核心問題：不靠任何醫療紀錄、也不直接量 BMI，只用日常生活行為，到底能不能準確預測一個人的肥胖風險？而且——最少需要幾個變數就夠？",
       "這就是我們整個專案要回答的問題。"], s1)


def s2(s):
    header(s, [("動機：從「落後指標 BMI」走向「", INK, True), ("領先指標 行為", BLUE, True), ("」", INK, True)])
    card(s, Inches(0.6), Inches(1.8), Inches(7.6), Inches(1.25), "🌍 全球肥胖危機",
         "全球逾 10 億人受肥胖困擾，台灣成人過重肥胖率超過 50%。傳統 BMI 是『事後測量』，無法提前預警。", BLUE)
    card(s, Inches(0.6), Inches(3.2), Inches(7.6), Inches(1.25), "📊 生活行為才是關鍵",
         "睡眠、吸菸、飲酒、運動是『領先指標』。機器學習可量化『行為 → 風險』的對應關係。", GREEN)
    card(s, Inches(0.6), Inches(4.6), Inches(7.6), Inches(1.25), "💡 本研究問題",
         "只用日常行為，能否準確預測肥胖風險？最少要幾個變數？", AMBER)
    add_img_fit(s, os.path.join(IMG, "m02_hero.png"), Inches(8.5), Inches(1.8), Inches(4.3), Inches(4.6))

build(2, "動機",
      [("落後指標 (lagging)", "結果發生後才看得到的指標，如 BMI——胖了才知道"),
       ("領先指標 (leading)", "結果發生前就能觀察、可提早介入的指標，如生活行為")],
      ["先講動機。肥胖是全球公衛問題，全球超過十億人、台灣成人過重肥胖率破五成。",
       "但傳統 BMI 是『落後指標』——你已經胖了才量得出來，沒辦法提前預警。",
       "我們的想法是：睡眠、吸菸、飲酒、運動這些『行為』是『領先指標』，能不能用它們提前預測風險？這就帶到我們的核心問題：只用行為、最少幾個變數就能預測。"], s2)


def s3(s):
    header(s, [("四版本迭代：", INK, True), ("從失敗中換來的判斷力", BLUE, True)],
           "每一次踩坑，都讓我們更清楚什麼才是真正可用的資料")
    rows = [("v1", "Lifestyle and Wellbeing (Kaggle)", "15,972 × 24", "❌ 資料不足", MUTE),
            ("v2", "Obesity Risk Dataset (Kaggle S4E2)", "20,800 × 17", "❌ 合成放大、BMI 不符 WHO", MUTE),
            ("v3", "Heart Disease heart_2020 (Kaggle)", "319,795 × 18", "❌ 第三方清洗18欄、2020疫情、無睡眠", MUTE),
            ("v4 ✅", "CDC BRFSS 2022 原始 XPT", "445k → 329k", "✅ 官方原始、328變數、含睡眠", GREEN)]
    y = Inches(2.0)
    for tag, name, size, why, col in rows:
        c = card(s, Inches(0.6), y, Inches(12.1), Inches(0.92),
                 [(tag + "  ", col, True), (name, INK, True)],
                 [[(size + "    ", MUTE, False), (why, col, False)]], col)
        y = y + Inches(1.06)

build(3, "四版本迭代",
      [("Kaggle", "全球最大的資料科學競賽與資料集分享平台"),
       ("合成放大資料", "用演算法(如 SMOTE)從少量真實樣本人工生成更多筆，分布可能失真"),
       ("XPT (.xpt)", "SAS 統計軟體的原始資料格式，CDC 官方釋出的原始檔")],
      ["我們的資料其實換了四個版本，這是踩坑換來的判斷力。",
       "v1 生活福祉資料太小、v2 是 Kaggle 競賽的合成放大資料、BMI 邊界不符 WHO；v3 是別人清洗過的 BRFSS 2020，只剩 18 欄、又是疫情年、還沒有睡眠變數。",
       "直到 v4——我們回到 CDC 官方原始 XPT 檔自己做 ETL，才真正可用。這也帶出我們最大的心得之一：回歸一手資料。",
       "（提醒：之後口頭都稱 v4 為『最終版』）"], s3)


def s4(s):
    header(s, [("資料來源：", INK, True), ("CDC BRFSS 2022", BLUE, True)],
           "全球最大的持續性健康行為電話調查 · 美國 CDC 主持")
    card(s, Inches(0.6), Inches(1.9), Inches(7.6), Inches(1.35), "🏛️ 什麼是 CDC",
         "美國疾病管制與預防中心，美國最高公共衛生機構，數據具權威性與全國代表性。", BLUE)
    card(s, Inches(0.6), Inches(3.4), Inches(7.6), Inches(1.5), "📞 什麼是 BRFSS",
         "行為風險因子監測系統。CDC 自 1984 年起每年電話訪問全美 50 州成人，年逾 40 萬人。", BLUE)
    card(s, Inches(0.6), Inches(5.05), Inches(3.7), Inches(1.7), "📦 原始規格",
         "SAS .XPT，445,132 筆 × 328 變數；2022 是最新含睡眠的年份。", GREEN)
    card(s, Inches(4.5), Inches(5.05), Inches(3.7), Inches(1.7), "🎯 為何用原始檔",
         "掌握 328 變數才能自由設計三層特徵集。", GREEN)
    add_img_fit(s, os.path.join(IMG, "s04_hero.png"), Inches(8.5), Inches(1.9), Inches(4.3), Inches(4.8))

build(4, "資料來源 CDC BRFSS 2022",
      [("CDC", "美國疾病管制與預防中心 (Centers for Disease Control and Prevention)"),
       ("BRFSS", "行為風險因子監測系統，CDC 自 1984 年起的年度健康電話調查"),
       ("SLEPTIM1", "BRFSS 中的睡眠時數欄位，偶數年沒有收集")],
      ["我們最終用的是 CDC 的 BRFSS 2022。CDC 是美國最高公衛機構，BRFSS 是它從 1984 年做到現在的年度健康電話調查，每年訪問超過 40 萬名美國成人，是全球最大的持續性健康監測。",
       "原始檔是 SAS 的 XPT 格式，44.5 萬筆、328 個變數。我們選 2022 的關鍵原因是：偶數年沒有睡眠時數，2022 是最新一年有睡眠的。",
       "用原始檔的好處是：抽樣有全國代表性，而且我們能掌握全部 328 個變數，才能自由設計後面的三層特徵集。"], s4)


def s5(s):
    header(s, [("ETL 清理流程：", INK, True), ("445,132 → 329,126 筆", BLUE, True)], "保留率 73.9%")
    steps = [("445,132", "原始 BRFSS 2022 (XPT)", 1.0, BLUE),
             ("396,326", "刪 BMI 缺值  −48,806", 0.89, BLUE),
             ("392,216", "移除睡眠無效值 (77/99)  −4,110", 0.88, BLUE),
             ("338,666", "刪核心欄位無效值  −53,550", 0.76, BLUE),
             ("329,126", "✅ 清理新增變數 + 修正 99900 異常碼 → 最終 26 欄", 0.74, GREEN)]
    y = Inches(2.0)
    for val, desc, frac, col in steps:
        bar = s.shapes.add_shape(1, Inches(0.6), y, Inches(0.6 + 9.0 * frac), Inches(0.62))
        bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
        btf = bar.text_frame; btf.margin_top = Inches(0.05); btf.margin_left = Inches(0.15)
        para(btf, val, 15, RGBColor(0x06,0x12,0x2B), bold=True, first=True, space=0)
        lb, ltf = tb(s, Inches(10.0), y, Inches(3.1), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
        para(ltf, desc, 11.5, MUTE, first=True, space=0)
        y = y + Inches(0.92)
    box, tf = tb(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.6))
    para(tf, "體重過輕重編為正常（n=6,778，改標籤不刪行）；每一步都基於資料字典逐欄判斷。", 13, MUTE, first=True)

build(5, "ETL 清理流程",
      [("ETL", "Extract-Transform-Load，資料萃取→清理轉換→載入的標準流程"),
       ("缺值 / 無效值", "問卷中『不知道/拒答』等代碼，不能當真實數值"),
       ("資料字典 (codebook)", "官方文件，說明每個欄位的編碼意義")],
      ["這頁是資料清理。原始 44.5 萬筆，我們一步步清掉缺值與無效值：刪 BMI 缺值、睡眠的無效代碼、核心欄位的無效值，最後清理我們新增的變數、修正一個叫 99900 的異常碼，剩下 32.9 萬筆、26 欄，保留率約 74%。",
       "每一步都不是亂刪，是對照 CDC 的資料字典逐欄判斷。例如體重過輕的人，我們是把標籤重編成正常、而不是刪掉，避免損失樣本。"], s5)


def s6(s):
    header(s, [("ETL 關鍵編碼決策", INK, True)], "魔鬼藏在編碼細節裡")
    cards = [("MENTHLTH / PHYSHLTH = 88", "心理／身體不適天數", "→ 0 天。BRFSS 用 88 代表『沒有任何一天不適』，不可當缺值刪（佔約 6 成樣本）。"),
             ("_DRNKWK2 = 99900", "每週估計飲酒量", "→ NaN。最終版新發現的異常編碼，實為無效值佔位（2,114 筆）。"),
             ("GENHLTH 反向處理", "自評整體健康", "原始 1=最好、5=最差，反向後讓數值方向與風險一致。"),
             ("連續 vs 類別", "特徵型別編碼策略", "連續特徵保留原值；類別特徵用 Label Encoding 轉整數碼。")]
    y = Inches(1.9)
    for hd, zh, body in cards:
        card(s, Inches(0.6), y, Inches(12.1), Inches(1.12),
             [(hd + "   ", INK, True), (zh, BLUE, False)], body, AMBER)
        y = y + Inches(1.22)

build(6, "ETL 關鍵編碼決策",
      [("Label Encoding", "把類別文字轉成整數代碼（如 男=1 女=0），樹模型可直接吃"),
       ("One-Hot Encoding", "把每個類別展成一個 0/1 欄位，會增加維度"),
       ("NaN", "Not a Number，代表缺失值")],
      ["這頁是幾個關鍵的編碼決策，也是最容易踩坑的地方。",
       "第一，心理和身體不適天數裡的 88，不是 88 天，而是『一天都沒有不適』，要轉成 0。如果當缺值刪，會一次錯刪六成樣本。",
       "第二，每週飲酒量有個 99900 的異常碼，其實是無效值佔位，要轉成缺失。第三，自評健康原始是 1 最好 5 最差，我們反向，讓數字越大代表越不健康、跟風險方向一致。",
       "最後，連續變數保留原值，類別變數用 Label Encoding。因為 XGBoost、隨機森林對這種編碼不敏感，不需要 One-Hot、避免維度爆炸。"], s6)


def s7(s):
    header(s, [("特徵集設計：", INK, True), ("三層階梯式消除", BLUE, True)],
           "嚴格包含關係 Core_6 ⊂ Behav_15 ⊂ Full_25，才能乾淨歸因準確率變化來自哪一類變數")
    card(s, Inches(0.6), Inches(2.1), Inches(6.6), Inches(1.2), [("Core_6", AMBER, True), ("（最簡可行模型）", MUTE, False)],
         "純可控行為：睡眠／吸菸／飲酒／運動／自評健康（6）", AMBER)
    card(s, Inches(0.6), Inches(3.5), Inches(6.6), Inches(1.2), [("Behav_15", GREEN, True), ("（移除慢性病）", MUTE, False)],
         "行為 + 健康狀態 + 人口統計（不含 7 項慢性病，15）", GREEN)
    card(s, Inches(0.6), Inches(4.9), Inches(6.6), Inches(1.2), [("Full_25", BLUE, True), ("（上限 Baseline）", MUTE, False)],
         "行為 + 健康 + 7 項慢性病史 + 人口統計（25）", BLUE)
    add_img_fit(s, os.path.join(IMG, "s07_pyramid.png"), Inches(7.6), Inches(2.0), Inches(5.2), Inches(4.6))

build(7, "特徵集設計：三層階梯",
      [("特徵集 (feature set)", "一組拿來訓練模型的輸入欄位組合"),
       ("⊂ 包含關係", "Core_6 完全被包含在 Behav_15 裡，Behav_15 又被 Full_25 包含"),
       ("Baseline", "基準線/對照組，這裡指『用最多變數』的上限版本")],
      ["這是我們方法的核心設計：三層階梯式的特徵集，而且是嚴格包含關係。",
       "最小的 Core_6 只有 6 個純可控行為；Behav_15 再加上健康狀態和人口統計，但刻意不放慢性病；最大的 Full_25 才把 7 項慢性病史也加進來。",
       "為什麼要嚴格包含？因為這樣當準確率改變時，我們能乾淨地歸因『是哪一類變數帶來的差異』——這是整個實驗能比較的基礎。"], s7)


def s8(s):
    header(s, [("兩種分類框架：", INK, True), ("二元 vs 三類別", BLUE, True)],
           "二元＝Normal vs 過重+肥胖；三類別＝Normal / Overweight / Obese")
    data = [("Full_25", "0.6867", "0.4955"), ("Behav_15", "0.6800", "0.4831"), ("Core_6", "0.6082", "0.4187")]
    y = Inches(2.0)
    for name, b, t in data:
        card(s, Inches(0.6), y, Inches(6.4), Inches(1.25),
             name, [[("二元 Binary F1  ", GREEN, True), (b, GREEN, True), ("    三類別 3-Class F1  ", MUTE, False), (t, MUTE, False)]], BLUE)
        y = y + Inches(1.35)
    box, tf = tb(s, Inches(0.6), Inches(6.2), Inches(6.5), Inches(1))
    para(tf, "二元普遍比三類別高 ~0.19 → 瓶頸是 Normal/Overweight 邊界模糊，非特徵不足。後續成績皆以二元為準。", 13, INK, first=True)
    add_img_fit(s, os.path.join(IMG, "binary_vs_3class.png"), Inches(7.4), Inches(1.9), Inches(5.4), Inches(5.0))

build(8, "兩種分類框架：二元 vs 三類別",
      [("二元分類 (Binary)", "把目標切成兩類：正常 vs 過重+肥胖"),
       ("三類別分類 (3-Class)", "切成三類：正常 / 過重 / 肥胖"),
       ("F1 分數", "精確率與召回率的調和平均，0~1，越高越好；不平衡資料常用")],
      ["在看成績前，要先講我們用兩種分類框架。二元是把人分成『正常』跟『過重+肥胖』兩類；三類別是分成正常、過重、肥胖三類。",
       "右圖很清楚：同一批模型，二元的 F1 是 0.687，三類別只有 0.495，差了快 0.19。原因是『過重』跟『正常』的 BMI 邊界本來就模糊，模型容易誤分——這是資料本質問題，不是特徵不夠。",
       "所以我們後面所有成績，都以二元為準。"], s8)


def s9(s):
    header(s, [("評估方法與指標", INK, True)], "公平比較 27 組的共同標準")
    card(s, Inches(0.6), Inches(1.9), Inches(3.9), Inches(2.4), "✂️ 資料切分",
         "80/20 train-test split，採 stratified 分層抽樣，確保各 BMI 類別比例在訓練集與測試集中一致。", BLUE)
    card(s, Inches(4.7), Inches(1.9), Inches(3.9), Inches(2.4), "📈 主指標：Macro F1",
         "三類不平衡下 Accuracy 易受多數類主導；Macro F1 對少數類更公平。", GREEN)
    card(s, Inches(8.8), Inches(1.9), Inches(3.9), Inches(2.4), "⚖️ 類別不平衡處理",
         "資料約 1:2（Normal 33% / OW+Obese 67%）。SMOTE 對少數類內插合成（僅訓練集），避免洩漏。", ORANGE)
    card(s, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.0), "為何選 SMOTE 而非欠採樣",
         ["SMOTE（過採樣）：特徵空間插值合成新少數類樣本、僅用於訓練集、27 組中最穩定。",
          "欠採樣：直接丟棄多數類 → 不可逆資訊損失、性能不穩。",
          "🔑 關鍵原則：SMOTE 僅套用於訓練集，測試集保持原始分佈，確保評估真實反映泛化能力。"], GREEN)

build(9, "評估方法與指標",
      [("train-test split", "把資料切成訓練集(學)與測試集(考)，避免用考過的題評分"),
       ("stratified 分層抽樣", "切分時維持各類別比例一致，避免某類被切偏"),
       ("Macro F1", "各類別 F1 取平均（不加權），小類別跟大類別一樣重要"),
       ("SMOTE", "Synthetic Minority Over-sampling，對少數類在特徵空間內插生成合成樣本"),
       ("資訊洩漏 (leakage)", "測試集資訊不慎進入訓練，導致成績虛高")],
      ["27 組模型要能公平比較，得有共同標準。三件事：",
       "一、資料切分用 80/20，並做分層抽樣，確保訓練跟測試的 BMI 比例一致。",
       "二、主指標用 Macro F1 不用 Accuracy，因為資料 1 比 2 不平衡，Accuracy 會被多數類灌水；Macro F1 對少數類比較公平。",
       "三、不平衡我們用 SMOTE 過採樣，而不是欠採樣——欠採樣直接丟資料、不可逆；SMOTE 是在特徵空間插值生成。最關鍵的原則是：SMOTE 只在訓練集做，測試集保持原始分佈，這樣評估才真實。"], s9)


def s10(s):
    header(s, [("模型訓練：", INK, True), ("27 組實驗矩陣", BLUE, True)], "3 特徵集 × 3 演算法 × 3 抽樣")
    # 簡化矩陣文字
    card(s, Inches(0.6), Inches(2.0), Inches(7.4), Inches(4.2), "實驗設計",
         ["演算法：LR (Logistic Regression) / RF (Random Forest) / XGBoost",
          "抽樣：Raw 原始 / SMOTE 過採樣 / Under-sampling 欠採樣",
          "特徵集：Core_6 / Behav_15 / Full_25",
          "→ 3 × 3 × 3 = 27 組模型，公平比較",
          "",
          "🏆 XGBoost + SMOTE 在三個特徵集全數奪冠。"], BLUE)
    add_img_fit(s, os.path.join(PLOT, "Full_25feat_f1_grid.png"), Inches(8.3), Inches(2.0), Inches(4.5), Inches(4.2))

build(10, "27 組實驗矩陣",
      [("LR", "邏輯迴歸，線性基準模型"),
       ("RF", "隨機森林，多棵決策樹投票，抗噪較強"),
       ("XGBoost", "梯度提升樹，表格資料最強的模型之一"),
       ("Raw / SMOTE / Under", "不處理 / 過採樣補少數類 / 欠採樣丟多數類")],
      ["訓練的部分，我們做了一個 27 組的實驗矩陣：3 個特徵集、乘以 3 種演算法、乘以 3 種抽樣策略。",
       "演算法是邏輯迴歸、隨機森林、XGBoost；抽樣是原始、SMOTE、欠採樣。這樣 3 乘 3 乘 3 等於 27 組，全部用同一套標準比較。",
       "結論很乾淨：XGBoost 加 SMOTE 這個組合，在三個特徵集全部拿第一。右邊熱圖就是其中一個特徵集的成績全貌，右下角最深就是這個最佳組合。"], s10)


def s11(s):
    header(s, [("結果：F1 成績總表", INK, True), ("（二元 Binary）", BLUE, True)], "XGBoost + SMOTE 三個特徵集均達最佳")
    bars = [("Full_25", "0.6867", 1.00, BLUE), ("Behav_15", "0.6800", 0.99, GREEN), ("Core_6", "0.6082", 0.886, AMBER)]
    bx = Inches(1.2); maxh = 2.9
    for name, val, frac, col in bars:
        h = maxh * frac
        top = Inches(1.9) + Inches(maxh - h)
        bar = s.shapes.add_shape(1, bx, top, Inches(2.0), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background(); bar.shadow.inherit = False
        vb, vtf = tb(s, bx, Inches(1.5), Inches(2.0), Inches(0.4))
        para(vtf, val, 20, INK, bold=True, first=True, align=PP_ALIGN.CENTER, space=0)
        cb, ctf = tb(s, bx, Inches(4.95), Inches(2.0), Inches(0.5))
        para(ctf, name, 14, INK, bold=True, first=True, align=PP_ALIGN.CENTER, space=0)
        bx = bx + Inches(2.5)
    card(s, Inches(8.3), Inches(1.9), Inches(4.4), Inches(2.1), "📋 三特徵集組成",
         ["Core_6：純行為 6", "Behav_15：+健康+人口", "Full_25：再+7 項慢性病史"], BLUE)
    card(s, Inches(8.3), Inches(4.2), Inches(4.4), Inches(2.4), "🏆 為何 XGBoost+SMOTE 最佳",
         ["XGBoost：擅長非線性與特徵交互", "SMOTE：平衡 1:2 類別", "→ 27 組三特徵集全數奪冠"], GREEN)

build(11, "F1 成績總表（二元）",
      [("過擬合 (overfitting)", "模型死背訓練資料、對新資料變差"),
       ("特徵交互", "兩個變數合起來才有意義的效果，樹模型擅長捕捉")],
      ["這是最終成績總表，都是二元 F1。Full_25 用 25 個變數是 0.687，Behav_15 用 15 個行為是 0.680，Core_6 只用 6 個是 0.608。",
       "右邊兩張卡解釋：左邊是三個特徵集各自的組成；右邊解釋為什麼 XGBoost 加 SMOTE 最好——XGBoost 擅長抓非線性跟特徵交互，SMOTE 平衡了 1 比 2 的類別，所以 27 組裡三個特徵集全勝。"], s11)


def s12(s):
    header(s, [("🔑 關鍵發現：", INK, True), ("行為勝過慢性病史", BLUE, True)], "數字皆為二元 Binary F1")
    card(s, Inches(0.8), Inches(2.3), Inches(3.4), Inches(1.8), "Full_25（含 7 項病歷）", "0.6867", BLUE)
    card(s, Inches(4.7), Inches(2.3), Inches(3.4), Inches(1.8), "Behav_15（無病歷）", "0.6800", GREEN)
    card(s, Inches(8.6), Inches(2.3), Inches(3.9), Inches(1.8), "差距", "< 1%（0.0067）", AMBER)
    box, tf = tb(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2))
    para(tf, "慢性病史的邊際貢獻 < 1%；Core_6 只問 6 題就達 F1=0.608，與 25 變數差距僅 0.078。", 17, INK, bold=True, first=True, space=10)
    para(tf, "→ 公衛應優先投資『行為介入』（睡眠、戒菸、運動），而非侵入式篩檢。", 16, GREEN, bold=True)

build(12, "關鍵發現：行為勝過慢性病史",
      [("邊際貢獻", "多加一組變數後，成績額外增加的幅度"),
       ("行為介入", "透過改變生活習慣來降低風險的公衛策略")],
      ["這是我們最重要的發現。比較 Full_25（含 7 項慢性病史）跟 Behav_15（完全沒有病史），F1 只差不到 1%。",
       "也就是說，慢性病史對預測肥胖的『邊際貢獻』不到 1%。而最精簡的 Core_6 只問 6 題，就有 0.608，跟用滿 25 個變數也只差 0.078。",
       "公衛上的意義很大：與其做侵入式的病史篩檢，不如優先投資行為介入——睡眠、戒菸、運動。"], s12)


def s13(s):
    header(s, [("特徵重要度排行", INK, True), ("（二元 Binary XGBoost）", BLUE, True)],
           "數值＝各特徵在所有樹分裂中的貢獻佔比（三組各自總和=1）")
    cols = [("Full_25", [("關節炎", "0.121"), ("運動", "0.111"), ("氣喘", "0.093"), ("性別", "0.092"), ("自評健康", "0.089"), ("看牙醫", "0.075")], BLUE),
            ("Behav_15", [("運動", "0.152"), ("性別", "0.146"), ("自評健康", "0.144"), ("行走困難", "0.120"), ("憂鬱", "0.107"), ("教育", "0.059")], GREEN),
            ("Core_6", [("自評健康", "0.418"), ("運動", "0.258"), ("吸菸", "0.120"), ("每週飲酒量", "0.120"), ("睡眠", "0.059"), ("是否飲酒", "0.026")], AMBER)]
    x = Inches(0.6)
    for name, items, col in cols:
        body = [f"{i+1}. {t}  {v}" for i, (t, v) in enumerate(items)]
        card(s, x, Inches(1.9), Inches(3.95), Inches(4.0), name, body, col)
        x = x + Inches(4.05)
    box, tf = tb(s, Inches(0.6), Inches(6.1), Inches(12), Inches(1))
    para(tf, "去掉慢性病後（Behav_15），運動/自評健康升到最前；Core_6 中自評健康+運動近 68%。", 13, MUTE, first=True)

build(13, "特徵重要度",
      [("特徵重要度", "模型認為每個變數對預測有多重要（XGBoost 用分裂增益計算）"),
       ("代理變數 (proxy)", "本身不是病因，但能間接反映其他狀態，如『看牙醫』反映健康意識")],
      ["這頁可以切換看三個特徵集的特徵重要度，都是真實模型輸出。",
       "重點是這個故事：在 Full_25 裡關節炎排第一，但前幾名已經有運動、自評健康、看牙醫這些行為跟意識指標。一旦把慢性病拿掉（Behav_15），運動、自評健康、行走困難立刻升到最前面。",
       "而 Core_6 裡，光自評健康就佔 0.42、運動 0.26，兩個可改變的行為合計快 7 成。這再次印證：行為才是主力。"], s13)


def s14(s):
    header(s, [("🎯 現場互動預測 Demo", INK, True)], "輸入你的生活習慣，立刻查看個人化肥胖風險")
    steps = [("01", "掃 QR 開啟網頁", "手機瀏覽器直接運行"),
             ("02", "選擇版本", "簡易版 6 項 / 完整版 15 項"),
             ("03", "填寫生活習慣", "睡眠、運動、飲酒等"),
             ("04", "查看風險儀表板", "🟢<35% 🟡35–65% 🔴>65%")]
    y = Inches(2.0)
    for no, t, d in steps:
        card(s, Inches(0.6), y, Inches(7.4), Inches(1.0), [(no + "  ", BLUE, True), (t, INK, True)], d, BLUE)
        y = y + Inches(1.12)
    add_img_fit(s, os.path.join(IMG, "s15_hero.png"), Inches(8.4), Inches(2.0), Inches(4.4), Inches(4.5))

build(14, "現場互動 Demo（步驟）",
      [("API", "應用程式介面，這裡指我們部署在雲端、可被網頁呼叫的預測服務"),
       ("儀表板 (dashboard)", "把預測結果用圖像化方式呈現的畫面")],
      ["接下來是現場 Demo。流程很簡單：掃 QR、選版本（6 項或 15 項）、填生活習慣、看風險儀表板。",
       "風險用顏色分級：綠色低於 35% 是低風險、黃色 35 到 65% 是邊界、紅色高於 65% 是高風險。等一下大家可以自己掃來玩。"], s14)


def s15(s):
    header(s, [("🎯 現場 Demo — ", INK, True), ("掃 QR 試打 ＋ 週追蹤", BLUE, True)])
    add_img_fit(s, os.path.join(IMG, "qr_demo.png"), Inches(1.3), Inches(2.0), Inches(3.6), Inches(3.6))
    box, tf = tb(s, Inches(0.9), Inches(5.7), Inches(4.4), Inches(0.8))
    para(tf, "① 同學現在掃 ↑ 開手機版填自己的習慣", 14, INK, bold=True, first=True, align=PP_ALIGN.CENTER)
    card(s, Inches(6.0), Inches(2.2), Inches(6.6), Inches(1.6), "▶ 單日 API",
         "開單日試打頁：填一天 → 真打一次 API 拿風險機率。", BLUE)
    card(s, Inches(6.0), Inches(4.0), Inches(6.6), Inches(1.6), "📈 開啟 16 週追蹤 Demo",
         "選 6 場景之一 → 逐週真打 16 次 API → 雙 Y 軸曲線（機率 + 模擬 BMI）。", AMBER)

build(15, "現場 Demo（操作）",
      [("FastAPI", "我們用來架設預測服務的 Python 框架"),
       (".pkl 模型檔", "訓練好的模型存成檔案，伺服器載入後即可即時推論")],
      ["這頁是台上實際操作。左邊大 QR 給同學掃，進到手機版填自己的習慣、真打一次 API 看風險。",
       "右邊兩顆是台上版：『單日 API』示範填一天、真的呼叫一次雲端 API 拿結果；『16 週追蹤』則是選一個情境、逐週呼叫 16 次 API，畫出機率跟模擬 BMI 的雙軸曲線。整條流程是手機填表→雲端 FastAPI 載入 pkl 模型推論→回傳機率。"], s15)


def s16(s):
    header(s, [("動態追蹤：", INK, True), ("改善 vs 惡化 × 16 週（Core_6）", BLUE, True)],
           "實線＝趨勢；圓點＝模型逐週真實輸出（即 API /predict/core6 結果）")
    card(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(2.0), "🟢 改善情境",
         ["起 58% → 末 32%（淨 −26）", "逐週：58→37→40→63→68→…→32", "區間 4.7% – 67.6%"], GREEN)
    card(s, Inches(6.8), Inches(2.1), Inches(5.9), Inches(2.0), "🔴 惡化情境",
         ["起 32% → 末 38%（淨 +6）", "逐週：32→28→5→5→40→62→…→38", "區間 4.7% – 67.6%"], RGBColor(0xEF,0x44,0x44))
    box, tf = tb(s, Inches(0.6), Inches(4.5), Inches(12.1), Inches(2))
    para(tf, "逐週在 4.7%–67.6% 劇烈震盪，淨變化卻僅 −26 / +6。", 17, INK, bold=True, first=True, space=10)
    para(tf, "→ 截面模型『無時間概念』，逐週是雜訊、只有起訖趨勢可信。這正是下一頁 Discussion 的限制。", 15, AMBER, bold=True)

build(16, "動態追蹤：改善 vs 惡化",
      [("截面資料 (cross-sectional)", "某一時間點對不同人的快照，沒有同一人隨時間的紀錄"),
       ("縱貫資料 (longitudinal)", "同一個人在多個時間點的追蹤紀錄")],
      ["這頁是動態追蹤。我們把『持續改善』跟『持續惡化』兩種行為序列，逐週丟進模型，這些值就是真實 API 輸出。",
       "但你會看到一個重點：圓點上下劇烈跳動，改善情境逐週在 4.7% 到 67.6% 之間亂跳，可是頭尾淨變化只有 26%。",
       "為什麼？因為我們的資料是『截面』的——它是不同人的快照、沒有同一個人的時間軸，所以模型逐週預測其實是雜訊，只有起訖趨勢可信。這就帶到下一頁的討論。"], s16)


def s17(s):
    header(s, [("討論：", INK, True), ("截面資料的根本限制", BLUE, True)])
    card(s, Inches(0.6), Inches(2.0), Inches(7.4), Inches(1.8), "觀察到的異常",
         "男性行為全部改善時，Core_6 機率正確下降，但 Behav_15 卻從 50% 升到 72%。", AMBER)
    card(s, Inches(0.6), Inches(4.0), Inches(7.4), Inches(2.0), "原因解析",
         "sex=1（男）+ age_group=6（中年）是『強錨點』；截面非縱向，模型學到『中年男性群體風險偏高』，無法捕捉個人縱向改變。", ORANGE)
    add_img_fit(s, os.path.join(IMG, "s18_diagram.png"), Inches(8.3), Inches(2.0), Inches(4.5), Inches(4.0))
    box, tf = tb(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.8))
    para(tf, "Core_6 無人口錨定，方向反而更穩定。", 15, GREEN, bold=True, first=True)

build(17, "討論：截面資料的限制",
      [("錨點 (anchor)", "模型過度依賴的強特徵，如性別、年齡，會壓過行為的影響"),
       ("人口統計變數", "性別、年齡、教育、收入等描述『你是誰』的變數")],
      ["討論這頁，講一個我們觀察到的異常。當一個男性把行為全部改善，Core_6 的風險正確下降；但 Behav_15 反而從 50% 升到 72%。",
       "原因是 Behav_15 含了性別、年齡這些人口變數，『中年男性』變成強錨點——模型學到『中年男性群體風險本來就高』，於是壓過了個人行為的改善。",
       "這恰恰說明截面資料的限制：它能比較『群體』，卻抓不到『同一個人』的縱向改變。反而沒有人口錨定的 Core_6，方向更穩定。"], s17)


def s18(s):
    header(s, [("結論：", INK, True), ("三個關鍵洞察", BLUE, True)])
    items = [("洞察一", "行為即足", "僅 6 個日常行為 F1 達 0.61；15 項達 0.68，與 25 變數差 <1%。預測肥胖不必碰醫療隱私。", "insight1.png"),
             ("洞察二", "資料為王", "換過 4 個資料集，前 3 個全失敗；回 CDC 原始 XPT 自行 ETL 才成功。模型再強也救不了壞資料。", "insight2.jpg"),
             ("洞察三", "準而不續", "能準確判斷『當下風險』，但逐週預測在 4.7%–67.6% 劇烈跳動。長期追蹤需縱貫／穿戴資料。", "insight3.jpg")]
    x = Inches(0.6)
    for tag, h, d, img in items:
        add_img_fit(s, os.path.join(IMG, img), x, Inches(1.7), Inches(3.95), Inches(1.9))
        c = card(s, x, Inches(3.75), Inches(3.95), Inches(3.0), [(tag + "  ", ORANGE, True), (h, ORANGE, True)], d, ORANGE)
        x = x + Inches(4.05)

build(18, "結論：三個關鍵洞察",
      [],
      ["最後總結三個關鍵洞察。",
       "洞察一『行為即足』：只用 6 個行為就有 0.61，15 項 0.68、跟全特徵差不到 1%，代表預測肥胖不必碰醫療隱私。",
       "洞察二『資料為王』：我們換了四個資料集，前三個都失敗，直到回 CDC 一手原始檔自己清理才成功——模型再強也救不了壞資料。",
       "洞察三『準而不續』：模型能判斷當下風險，但逐週會劇烈跳動，個人化長期追蹤要靠縱貫或穿戴資料，這也是我們的未來工作。"], s18)


def s19(s):
    box, tf = tb(s, Inches(0), Inches(2.4), Inches(13.333), Inches(2.6), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "機器學習期末報告 · 第五組", 16, MUTE, first=True, align=PP_ALIGN.CENTER, space=16)
    para(tf, [("Thank You", INK, True), (".", BLUE, True)], 54, INK, bold=True, align=PP_ALIGN.CENTER, space=14)
    para(tf, "謝謝聆聽 — 歡迎提問 Q & A", 22, MUTE, align=PP_ALIGN.CENTER)

build(19, "結尾 Thank You",
      [],
      ["以上就是我們的報告，謝謝大家聆聽，歡迎提問。"], s19)


def s20(s):
    header(s, [("資料與資源來源 · ", INK, True), ("References", BLUE, True)])
    card(s, Inches(0.6), Inches(1.9), Inches(3.95), Inches(4.6), "📂 資料集演進",
         ["v1 Lifestyle & Wellbeing (Kaggle ydalat)",
          "v2 Obesity Risk S4E2 (Kaggle jpkochar)",
          "v3 heart_2020 (Kaggle kamilpytlak)",
          "v4 ✅ CDC BRFSS 2022 原始 XPT",
          "cdc.gov/brfss"], BLUE)
    card(s, Inches(4.7), Inches(1.9), Inches(3.95), Inches(4.6), "📚 學術文獻",
         ["Jastreboff AM, et al. (2022) Tirzepatide. NEJM 387(3):205-216 (SURMOUNT-1, −20.9%)",
          "Wilding JPH, et al. (2021) Semaglutide. NEJM 384(11):989-1002 (STEP-1, −14.9%)"], ORANGE)
    card(s, Inches(8.8), Inches(1.9), Inches(3.95), Inches(4.6), "🔗 專案資源",
         ["Demo 網站 inorino.github.io/obesity-demo",
          "API + Swagger obesity-api-2-2ro4.onrender.com/docs",
          "GitHub github.com/inorino/obesity-api-2"], GREEN)

build(20, "References 來源",
      [],
      ["這是我們的資料來源與參考文獻：四個資料集的演進、兩篇減重藥物的 NEJM 文獻，以及 Demo 網站、API、GitHub 連結。歡迎掃 QR 或上 GitHub 看我們的程式碼。"], s20)


prs.save(OUT_PPTX)
with open(OUT_MD, "w", encoding="utf-8") as f:
    f.write("\n".join(MD))
print("PPTX exists after save:", os.path.exists(OUT_PPTX), os.path.getsize(OUT_PPTX))
print("MD  :", os.path.exists(OUT_MD))

# ---- 轉 PDF（ascii 暫名避開 COM 中文路徑 bug）----
OUT_PDF = os.path.join(BASE, "outputs", "簡報_第五組.pdf")
try:
    import win32com.client
    tmp_pptx = os.path.join(BASE, "outputs", "_deck_tmp.pptx")
    tmp_pdf = os.path.join(BASE, "outputs", "_deck_tmp.pdf")
    shutil_copy = __import__("shutil").copyfile
    shutil_copy(OUT_PPTX, tmp_pptx)
    app = win32com.client.Dispatch("PowerPoint.Application")
    try:
        pres = app.Presentations.Open(tmp_pptx, ReadOnly=True, WithWindow=False)
        pres.SaveAs(tmp_pdf, 32)
        pres.Close()
    finally:
        app.Quit()
    if os.path.exists(OUT_PDF):
        os.remove(OUT_PDF)
    os.replace(tmp_pdf, OUT_PDF)
    os.remove(tmp_pptx)
    print("PDF :", os.path.exists(OUT_PDF), os.path.getsize(OUT_PDF))
except Exception as e:
    print("PDF FAILED:", repr(e))
