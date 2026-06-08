"""
gen_slides.py v2 — 深海Teal + 時間軸 + 進度條 + 標籤tag
"""
import sys
sys.stdout.reconfigure(encoding="utf-8")

# ─── 色盤（深海 Teal）────────────────────────────────────────
C = {
    "navy":   "#0a1628",
    "panel":  "#1e3a5f",
    "teal":   "#0d9488",
    "teal2":  "#14b8a6",
    "orange": "#f97316",
    "white":  "#ffffff",
    "lgray":  "#f1f5f9",
    "text":   "#1e293b",
    "muted":  "#64748b",
    "green":  "#16a34a",
    "red":    "#dc2626",
    "amber":  "#d97706",
}

# ─── 投影片內容 ───────────────────────────────────────────────
SLIDES = [
    {
        "type": "cover",
        "title": "生活習慣與肥胖風險預測",
        "subtitle": "CDC BRFSS 2022  ·  XGBoost  ·  Binary Classification",
        "meta": "329,126 筆有效樣本  ·  25 特徵  ·  Weighted F1 = 0.687",
        "tags": ["機器學習", "公共衛生", "BRFSS 2022"],
    },
    {
        "type": "timeline",
        "title": "研究歷程：三次迭代演進",
        "subtitle": "每次遇到資料問題 → 分析根因 → 換更好的資料源",
        "nodes": [
            {
                "label": "Iteration 01",
                "title": "Kaggle 合成資料",
                "status": "fail",
                "color": "red",
                "items": ["obesity_level.csv", "深度學習合成", "❌ BMI 邊界不符 WHO", "❌ 模型學假規律"],
            },
            {
                "label": "Iteration 02",
                "title": "BRFSS 2020 第三方",
                "status": "fail",
                "color": "amber",
                "items": ["heart_2020_cleaned.csv", "僅 18 欄位", "❌ COVID 第一年偏差", "❌ 缺睡眠變數"],
            },
            {
                "label": "Iteration 03",
                "title": "CDC BRFSS 2022 官方",
                "status": "success",
                "color": "teal",
                "items": ["官方 XPT 原始檔", "445,132 × 328 變數", "✅ 官方 BMI 公式", "✅ 完整行為特徵"],
            },
        ],
    },
    {
        "type": "section",
        "label": "01",
        "title": "資料  ETL",
        "subtitle": "收集、清理、轉換",
    },
    {
        "type": "split",
        "title": "BRFSS 2022 清理流程",
        "left": [
            ("01", "原始資料",  "445,132 筆 × 328 變數（CDC 官方 XPT 格式）"),
            ("02", "特徵選取",  "手動選取 25 個研究相關特徵"),
            ("03", "清理缺值",  "移除拒答碼 7/77、不知道碼 9/99"),
            ("04", "BMI 計算", "WHO 標準：體重(kg) ÷ 身高(m)²"),
            ("05", "標籤建立",  "Normal / Overweight / Obese"),
        ],
        "right": [("329,126", "有效樣本"), ("25", "特徵數"), ("26", "總欄位數")],
        "tags": ["XPT 格式", "WHO 標準", "CDC 官方"],
    },
    {
        "type": "cards3",
        "title": "三層階梯式特徵集設計",
        "subtitle": "核心設計：移除慢性病史，測試行為特徵的獨立預測力",
        "cards": [
            {
                "label": "Full_25feat", "num": "25", "unit": "特徵",
                "color": "teal",
                "items": ["行為習慣", "人口統計", "慢性病史"],
                "note": "完整基線（上限）",
            },
            {
                "label": "Behav_15feat", "num": "15", "unit": "特徵",
                "color": "orange", "highlight": True,
                "items": ["行為習慣", "人口統計", "移除慢性病史"],
                "note": "★ 無需醫療紀錄",
            },
            {
                "label": "Core_6feat", "num": "6", "unit": "特徵",
                "color": "panel",
                "items": ["睡眠 / 吸菸", "飲酒 / 運動", "自評健康"],
                "note": "最簡問卷",
            },
        ],
    },
    {
        "type": "section",
        "label": "02",
        "title": "演算法 / 方法",
        "subtitle": "Methodology",
    },
    {
        "type": "split",
        "title": "為什麼選 Binary Classification",
        "left": [
            ("❌", "3-class 嘗試",  "Normal / OW / Obese  →  F1 = 0.495"),
            ("🔍", "根因分析",      "Normal 與 OW 的 BMI 邊界極模糊，系統性誤判"),
            ("✅", "改為 Binary",   "Normal vs OW+Obese（BMI 25 為界）"),
            ("📈", "提升結果",      "F1 提升至 0.687（+0.192）"),
        ],
        "right": [("0.495", "3-class F1"), ("+0.19", "提升幅度"), ("0.687", "Binary F1")],
        "tags": ["Binary", "BMI 25 閾值", "Weighted F1"],
    },
    {
        "type": "bullets",
        "title": "實驗矩陣：27 個模型全面比較",
        "items": [
            "3 特徵集  ×  3 演算法  ×  3 抽樣策略  =  27 個模型",
            "演算法：Logistic Regression  /  Random Forest  /  XGBoost",
            "抽樣策略：Original  /  SMOTE  /  Random Undersampling",
            "資料不平衡：Normal 約 33%，Overweight+Obese 約 67%",
            "→ 最佳組合：XGBoost + SMOTE（三個特徵集一致）",
        ],
        "tags": ["XGBoost", "SMOTE", "27 Models"],
    },
    {
        "type": "bullets",
        "title": "類別不平衡與 SMOTE",
        "items": [
            "資料分佈：Normal 33%  vs  OW+Obese 67%，約 1:2 不平衡",
            "直接訓練：模型偏向多數類別，Normal 召回率偏低",
            "SMOTE：對 Normal 合成新樣本，使訓練集平衡",
            "Undersampling：丟棄多數類別資料，資訊損失較大",
            "→ SMOTE 在所有特徵集 × 演算法組合中表現最穩定",
        ],
        "tags": ["SMOTE", "類別不平衡", "資料增強"],
    },
    {
        "type": "section",
        "label": "03",
        "title": "比較與討論",
        "subtitle": "Comparison / Discussion",
    },
    {
        "type": "results",
        "title": "模型成績總覽（Binary Weighted F1）",
        "subtitle": "最佳演算法：XGBoost + SMOTE（三個特徵集一致）",
        "bars": [
            {"label": "Full_25feat",  "feat": "25 特徵", "value": 0.6867, "color": "teal",   "tag": "完整版"},
            {"label": "Behav_15feat", "feat": "15 特徵", "value": 0.6800, "color": "orange", "tag": "無醫療紀錄", "highlight": True},
            {"label": "Core_6feat",   "feat": " 6 特徵", "value": 0.6082, "color": "panel",  "tag": "純行為"},
        ],
    },
    {
        "type": "bullets",
        "title": "關鍵發現",
        "items": [
            "Behav_15 (0.680)  vs  Full_25 (0.687)：差距僅 0.007（< 1%）",
            "移除全部 7 個慢性病史特徵，F1 幾乎不變",
            "Core_6 (0.608)：只靠 6 個日常行為，仍達 60.8% F1",
            "瓶頸是 Normal/OW BMI 邊界模糊，非特徵數量問題",
            "→ 不需任何醫療紀錄，行為習慣即可有效預測肥胖風險",
        ],
        "tags": ["關鍵洞察", "行為介入", "公衛政策"],
    },
    {
        "type": "bullets",
        "title": "結論",
        "items": [
            "只靠 6 個行為特徵（Core_6）即可達到 F1 = 0.608",
            "加入人口統計擴展至 15 特徵（Behav_15），F1 = 0.680",
            "幾乎等同加入全部慢性病診斷的完整版（0.687），差距 < 1%",
            "支持以行為介入為核心的公衛政策，無需侵入式醫療資料",
            "→ 模型已部署為 REST API，支援前端即時查詢",
        ],
        "note": "API：https://obesity-api-2-2ro4.onrender.com",
        "tags": ["REST API", "行為介入", "公衛政策"],
    },
]


# ═══════════════════════════════════════════════════════════
# PPTX Generator
# ═══════════════════════════════════════════════════════════
def gen_pptx(slides, out):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height
    BLANK = prs.slide_layouts[6]

    def _rgb(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    def c(key):
        return _rgb(C[key])

    def bg(slide, key):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = c(key)

    def rect(slide, l, t, w, h, fill=None, fhex=None):
        tb = slide.shapes.add_textbox(l, t, w, h)
        if fill:
            tb.fill.solid()
            tb.fill.fore_color.rgb = c(fill)
        elif fhex:
            tb.fill.solid()
            tb.fill.fore_color.rgb = _rgb(fhex)
        tb.line.fill.background()
        return tb

    def oval(slide, l, t, w, h, fill="teal"):
        shp = slide.shapes.add_shape(9, l, t, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = c(fill)
        shp.line.fill.background()
        return shp

    def txt(slide, text, l, t, w, h,
            size=18, bold=False, fk="text", fhex=None,
            align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(fhex) if fhex else c(fk)
        run.font.name = "微軟正黑體"
        return tb

    def tag(slide, text, l, t, fill="teal"):
        tw = Inches(max(len(text) * 0.17 + 0.28, 0.9))
        th = Inches(0.28)
        tb = rect(slide, l, t, tw, th, fill=fill)
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = c("white")
        run.font.name = "微軟正黑體"
        return tw

    def tags_row(slide, tag_list, l, t):
        x = l
        for i, (text, fill) in enumerate(tag_list):
            w = tag(slide, text, x, t, fill=fill)
            x += w + Inches(0.1)

    def header(slide, title):
        rect(slide, 0, 0, W, Inches(0.07), fill="teal")
        txt(slide, title,
            Inches(0.55), Inches(0.1), W - Inches(1.1), Inches(0.75),
            size=23, bold=True, fk="navy")
        rect(slide, Inches(0.55), Inches(0.88),
             W - Inches(1.1), Inches(0.018), fhex="#cbd5e1")

    # ── Cover ─────────────────────────────────────────────
    def render_cover(slide, sl):
        bg(slide, "navy")
        rect(slide, 0, 0, Inches(0.1), H, fill="teal")
        rect(slide, 0, H - Inches(0.07), W, Inches(0.07), fill="orange")
        rect(slide, Inches(0.55), Inches(2.6), Inches(0.06), Inches(2.5), fill="orange")

        txt(slide, sl["title"],
            Inches(1.0), Inches(1.5), W - Inches(2.0), Inches(1.7),
            size=40, bold=True, fk="white")
        txt(slide, sl["subtitle"],
            Inches(1.0), Inches(3.3), W - Inches(2.0), Inches(0.6),
            size=19, fhex="#5eead4")
        txt(slide, sl.get("meta", ""),
            Inches(1.0), Inches(4.0), W - Inches(2.0), Inches(0.5),
            size=13, fk="muted")

        tag_fills = ["teal", "panel", "orange"]
        x = Inches(1.0)
        for i, t_text in enumerate(sl.get("tags", [])):
            w = tag(slide, t_text, x, Inches(4.8), fill=tag_fills[i % 3])
            x += w + Inches(0.15)

    # ── Timeline ──────────────────────────────────────────
    def render_timeline(slide, sl):
        bg(slide, "lgray")
        rect(slide, 0, 0, W, Inches(0.07), fill="teal")
        txt(slide, sl["title"],
            Inches(0.55), Inches(0.12), W - Inches(1.1), Inches(0.65),
            size=21, bold=True, fk="navy")
        txt(slide, sl.get("subtitle", ""),
            Inches(0.55), Inches(0.75), W - Inches(1.1), Inches(0.35),
            size=12, fk="muted", italic=True)

        nodes = sl["nodes"]
        n = len(nodes)
        line_y = Inches(2.35)
        xs = [Inches(1.8) + (W - Inches(3.6)) * i / (n - 1) for i in range(n)]

        # Horizontal connector line
        rect(slide, xs[0], line_y + Inches(0.18),
             xs[-1] - xs[0], Inches(0.04), fhex="#94a3b8")

        for i, (node, nx) in enumerate(zip(nodes, xs)):
            node_fill = node["color"]
            cr = Inches(0.4)

            # Circle
            circ = oval(slide, nx - cr, line_y, cr * 2, cr * 2, fill=node_fill)
            tf = circ.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(i + 1)
            run.font.size = Pt(17)
            run.font.bold = True
            run.font.color.rgb = c("white")
            run.font.name = "微軟正黑體"

            # Label above
            txt(slide, node["label"],
                nx - Inches(1.0), line_y - Inches(0.55),
                Inches(2.0), Inches(0.38),
                size=11, fk="muted", align=PP_ALIGN.CENTER)

            # Card below
            card_y = line_y + Inches(0.95)
            card_w = Inches(3.4)
            card_h = Inches(3.6)
            card_l = nx - card_w / 2

            is_ok = node["status"] == "success"
            card_fill = "navy" if is_ok else "white"
            rect(slide, card_l + Inches(0.05), card_y + Inches(0.05),
                 card_w, card_h, fhex="#94a3b8")  # shadow
            rect(slide, card_l, card_y, card_w, card_h, fill=card_fill)
            rect(slide, card_l, card_y, card_w, Inches(0.08), fill=node_fill)

            title_fk = "white" if is_ok else "navy"
            txt(slide, node["title"],
                card_l + Inches(0.2), card_y + Inches(0.15),
                card_w - Inches(0.4), Inches(0.48),
                size=14, bold=True, fk=title_fk)

            for j, item in enumerate(node["items"]):
                item_fk = "white" if is_ok else "text"
                txt(slide, item,
                    card_l + Inches(0.2),
                    card_y + Inches(0.72) + Inches(0.56) * j,
                    card_w - Inches(0.4), Inches(0.5),
                    size=12, fk=item_fk)

    # ── Section ───────────────────────────────────────────
    def render_section(slide, sl):
        bg(slide, "navy")
        rect(slide, 0, H * 0.42, W, Inches(0.05), fill="teal")
        txt(slide, sl["label"],
            Inches(1.5), Inches(1.3), Inches(3.5), Inches(1.5),
            size=80, bold=True, fhex=C["teal"])
        txt(slide, sl["title"],
            Inches(1.5), H * 0.46, W - Inches(3.0), Inches(1.2),
            size=36, bold=True, fk="white")
        if sl.get("subtitle"):
            txt(slide, sl["subtitle"],
                Inches(1.5), H * 0.46 + Inches(1.25), W - Inches(3.0), Inches(0.55),
                size=16, fhex="#5eead4")

    # ── Split layout ──────────────────────────────────────
    def render_split(slide, sl):
        bg(slide, "white")
        header(slide, sl["title"])

        left_items = sl["left"]
        for i, row in enumerate(left_items):
            num, title, desc = row
            y = Inches(1.02) + Inches(1.05) * i
            is_emoji = not str(num).replace("0","").isdigit()

            if is_emoji:
                txt(slide, num,
                    Inches(0.5), y + Inches(0.04), Inches(0.5), Inches(0.42),
                    size=20, align=PP_ALIGN.CENTER)
            else:
                circ = oval(slide, Inches(0.5), y + Inches(0.04),
                            Inches(0.42), Inches(0.42), fill="teal")
                tf = circ.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = num
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = c("white")
                run.font.name = "微軟正黑體"

            txt(slide, title,
                Inches(1.1), y, Inches(5.9), Inches(0.42),
                size=14, bold=True, fk="navy")
            txt(slide, desc,
                Inches(1.1), y + Inches(0.4), Inches(5.9), Inches(0.48),
                size=12, fk="muted")

        # Right panel
        pl = Inches(7.6)
        pw = W - pl - Inches(0.35)
        rect(slide, pl, Inches(0.93), pw, Inches(5.6), fill="navy")
        rect(slide, pl, Inches(0.93), Inches(0.06), Inches(5.6), fill="teal")

        right = sl["right"]
        ph = Inches(5.6) / len(right)
        for i, (val, label) in enumerate(right):
            cy = Inches(0.93) + ph * i
            txt(slide, val,
                pl + Inches(0.2), cy + Inches(0.25),
                pw - Inches(0.3), Inches(1.0),
                size=34, bold=True, fk="teal", align=PP_ALIGN.CENTER)
            txt(slide, label,
                pl + Inches(0.2), cy + ph - Inches(0.5),
                pw - Inches(0.3), Inches(0.4),
                size=12, fk="muted", align=PP_ALIGN.CENTER)

        if sl.get("tags"):
            fills = ["teal", "panel", "orange", "teal2"]
            x = Inches(0.55)
            for j, t_text in enumerate(sl["tags"]):
                w = tag(slide, t_text, x, H - Inches(0.52), fill=fills[j % 4])
                x += w + Inches(0.1)

    # ── Cards3 ────────────────────────────────────────────
    def render_cards3(slide, sl):
        bg(slide, "lgray")
        header(slide, sl["title"])
        if sl.get("subtitle"):
            txt(slide, sl["subtitle"],
                Inches(0.55), Inches(0.9), W - Inches(1.1), Inches(0.35),
                size=11, fk="muted", italic=True)

        cards = sl["cards"]
        gap = Inches(0.28)
        cw = (W - Inches(1.1) - gap * (len(cards) - 1)) / len(cards)
        ch = Inches(5.0)
        cy = Inches(1.35)

        for i, card in enumerate(cards):
            cl = Inches(0.55) + (cw + gap) * i
            is_hl = card.get("highlight", False)

            rect(slide, cl + Inches(0.05), cy + Inches(0.05),
                 cw, ch, fhex="#94a3b8")
            rect(slide, cl, cy, cw, ch, fill="white")
            rect(slide, cl, cy, cw, Inches(0.1), fill=card["color"])

            txt(slide, card["num"],
                cl, cy + Inches(0.15), cw, Inches(1.1),
                size=52, bold=True, fk=card["color"], align=PP_ALIGN.CENTER)
            txt(slide, card["unit"],
                cl, cy + Inches(1.2), cw, Inches(0.32),
                size=13, fk="muted", align=PP_ALIGN.CENTER)
            txt(slide, card["label"],
                cl + Inches(0.2), cy + Inches(1.6),
                cw - Inches(0.4), Inches(0.46),
                size=13, bold=True, fk="navy", align=PP_ALIGN.CENTER)

            for j, item in enumerate(card["items"]):
                txt(slide, "▸  " + item,
                    cl + Inches(0.2),
                    cy + Inches(2.2) + Inches(0.5) * j,
                    cw - Inches(0.4), Inches(0.44),
                    size=12, fk="text")

            note_y = cy + ch - Inches(0.52)
            nb = rect(slide, cl + Inches(0.2), note_y,
                      cw - Inches(0.4), Inches(0.36),
                      fill="teal" if is_hl else "muted")
            tf = nb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = card.get("note", "")
            run.font.size = Pt(11)
            run.font.bold = is_hl
            run.font.color.rgb = c("white")
            run.font.name = "微軟正黑體"

    # ── Bullets ───────────────────────────────────────────
    def render_bullets(slide, sl):
        bg(slide, "white")
        header(slide, sl["title"])

        y = Inches(1.02)
        for item in sl["items"]:
            is_arrow = str(item).startswith("→")
            tb = slide.shapes.add_textbox(
                Inches(0.55), y, W - Inches(1.1), Inches(0.62))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.space_before = Pt(3)
            run = p.add_run()
            if is_arrow:
                run.text = item
                run.font.bold = True
                run.font.color.rgb = c("teal")
            else:
                run.text = "▸  " + item
                run.font.color.rgb = c("text")
            run.font.size = Pt(16)
            run.font.name = "微軟正黑體"
            y += Inches(0.65)

        if sl.get("note"):
            txt(slide, sl["note"],
                Inches(0.55), H - Inches(0.55), W - Inches(1.1), Inches(0.38),
                size=11, fk="muted", italic=True)

        if sl.get("tags"):
            fills = ["teal", "panel", "orange", "teal2"]
            x = Inches(0.55)
            for j, t_text in enumerate(sl["tags"]):
                w = tag(slide, t_text, x, H - Inches(0.52), fill=fills[j % 4])
                x += w + Inches(0.1)

    # ── Results (progress bars) ───────────────────────────
    def render_results(slide, sl):
        bg(slide, "white")
        header(slide, sl["title"])
        if sl.get("subtitle"):
            txt(slide, sl["subtitle"],
                Inches(0.55), Inches(0.9), W - Inches(1.1), Inches(0.35),
                size=11, fk="muted", italic=True)

        bar_l   = Inches(3.8)
        bar_w   = Inches(7.2)
        bar_h   = Inches(0.52)
        bar_gap = Inches(1.1)
        start_y = Inches(1.55)

        for i, bar in enumerate(sl["bars"]):
            by = start_y + bar_gap * i
            is_hl = bar.get("highlight", False)

            if is_hl:
                rect(slide, Inches(0.3), by - Inches(0.14),
                     W - Inches(0.6), bar_h + Inches(0.28), fhex="#f0fdfa")

            txt(slide, bar["label"],
                Inches(0.55), by + Inches(0.08),
                Inches(2.5), Inches(0.38),
                size=14, bold=is_hl, fk="navy")

            # Feat tag
            ft = rect(slide, Inches(0.55), by + Inches(0.5),
                      Inches(1.1), Inches(0.26), fill=bar["color"])
            tf = ft.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = bar["feat"]
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = c("white")
            run.font.name = "微軟正黑體"

            # Bar bg
            rect(slide, bar_l, by + Inches(0.08),
                 bar_w, bar_h, fhex="#e2e8f0")

            # Bar fill
            fill_w = bar_w * bar["value"]
            rect(slide, bar_l, by + Inches(0.08),
                 fill_w, bar_h, fill=bar["color"])

            # Value
            txt(slide, f"  {bar['value']:.4f}",
                bar_l + fill_w, by + Inches(0.1),
                Inches(1.2), Inches(0.38),
                size=14, bold=True, fk="navy")

            # Right tag
            if bar.get("tag"):
                tag(slide, bar["tag"],
                    W - Inches(1.8), by + Inches(0.14),
                    fill=bar["color"])

        txt(slide,
            "▸ 三個特徵集最佳演算法一致：XGBoost + SMOTE",
            Inches(0.55), H - Inches(0.85),
            W - Inches(1.1), Inches(0.45),
            size=13, bold=True, fk="teal")

    # ── Render ────────────────────────────────────────────
    RENDERERS = {
        "cover":    render_cover,
        "timeline": render_timeline,
        "section":  render_section,
        "split":    render_split,
        "cards3":   render_cards3,
        "bullets":  render_bullets,
        "results":  render_results,
    }

    for sl in slides:
        slide = prs.slides.add_slide(BLANK)
        fn = RENDERERS.get(sl["type"])
        if fn:
            fn(slide, sl)

    prs.save(out)
    print(f"PPTX  →  {out}")


# ═══════════════════════════════════════════════════════════
# Reveal.js Generator（保持原有功能）
# ═══════════════════════════════════════════════════════════
def gen_reveal(slides, out):

    def esc(s):
        return str(s).replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")

    secs = []

    for sl in slides:
        t = sl["type"]

        if t == "cover":
            tag_html = "".join(
                f'<span style="background:{C["teal"]};color:#fff;'
                f'padding:3px 12px;border-radius:4px;font-size:0.65em;'
                f'margin-right:8px">{esc(x)}</span>'
                for x in sl.get("tags", [])
            )
            secs.append(f"""
<section data-background-color="{C['navy']}" data-transition="fade">
  <div style="padding:1em 2em;border-left:5px solid {C['teal']}">
    <h1 style="color:#fff;font-size:2.1em;font-weight:800;line-height:1.3;margin-bottom:0.4em">
      {esc(sl['title'])}</h1>
    <p style="color:#5eead4;font-size:1em;margin:0.3em 0">{esc(sl['subtitle'])}</p>
    <p style="color:{C['muted']};font-size:0.78em;margin-top:1em">{esc(sl.get('meta',''))}</p>
    <div style="margin-top:1.4em">{tag_html}</div>
  </div>
</section>""")

        elif t == "timeline":
            nodes_html = ""
            for i, node in enumerate(sl["nodes"]):
                clr = C.get(node["color"], C["teal"])
                is_ok = node["status"] == "success"
                bg_clr = C["navy"] if is_ok else "#fff"
                fg_clr = "#fff" if is_ok else C["text"]
                items_html = "".join(f'<li style="font-size:0.8em;margin:4px 0;color:{fg_clr}">{esc(x)}</li>' for x in node["items"])
                nodes_html += f"""
<div style="flex:1;text-align:center">
  <div style="font-size:0.65em;color:{C['muted']};margin-bottom:8px">{esc(node['label'])}</div>
  <div style="width:44px;height:44px;border-radius:50%;background:{clr};
              color:#fff;font-size:1.1em;font-weight:700;
              display:flex;align-items:center;justify-content:center;
              margin:0 auto 12px">{i+1}</div>
  <div style="background:{bg_clr};border-top:4px solid {clr};
              padding:14px;min-height:180px;text-align:left">
    <div style="font-weight:700;font-size:0.85em;color:{fg_clr if is_ok else C['navy']};margin-bottom:8px">
      {esc(node['title'])}</div>
    <ul style="list-style:none;margin:0;padding:0">{items_html}</ul>
  </div>
</div>"""
            secs.append(f"""
<section data-background-color="{C['lgray']}" data-transition="slide">
  <h3 style="color:{C['navy']};font-size:1.2em;font-weight:700;text-align:left;margin-bottom:0.3em">
    {esc(sl['title'])}</h3>
  <p style="color:{C['muted']};font-size:0.7em;text-align:left;margin-bottom:1em;font-style:italic">
    {esc(sl.get('subtitle',''))}</p>
  <div style="display:flex;gap:20px;align-items:flex-start">{nodes_html}</div>
</section>""")

        elif t == "section":
            secs.append(f"""
<section data-background-color="{C['navy']}" data-transition="slide">
  <div style="text-align:left;padding:0 2em">
    <p style="color:{C['teal']};font-size:4em;font-weight:900;line-height:1;margin:0 0 0.1em">
      {esc(sl['label'])}</p>
    <h2 style="color:#fff;font-size:2em;font-weight:700;margin:0">{esc(sl['title'])}</h2>
    <p style="color:#5eead4;font-size:0.85em;margin-top:0.5em">{esc(sl.get('subtitle',''))}</p>
    <div style="height:4px;background:{C['teal']};width:160px;margin-top:0.8em"></div>
  </div>
</section>""")

        elif t == "split":
            rows_html = ""
            for num, title, desc in sl["left"]:
                rows_html += f"""
<div class="fragment" style="display:flex;gap:12px;margin:10px 0;align-items:flex-start">
  <div style="min-width:32px;height:32px;border-radius:50%;
              background:{C['teal']};color:#fff;font-size:0.75em;
              font-weight:700;display:flex;align-items:center;
              justify-content:center;flex-shrink:0">{esc(num)}</div>
  <div>
    <div style="font-weight:700;font-size:0.82em;color:{C['navy']}">{esc(title)}</div>
    <div style="font-size:0.72em;color:{C['muted']}">{esc(desc)}</div>
  </div>
</div>"""
            right_html = "".join(
                f'<div style="text-align:center;padding:14px 0;border-bottom:1px solid #1e3a5f">'
                f'<div style="font-size:2em;font-weight:800;color:{C["teal"]}">{esc(v)}</div>'
                f'<div style="font-size:0.65em;color:{C["muted"]}">{esc(l)}</div></div>'
                for v, l in sl["right"]
            )
            tag_html = "".join(
                f'<span style="background:{C["teal"]};color:#fff;padding:2px 10px;'
                f'border-radius:4px;font-size:0.55em;margin-right:6px">{esc(x)}</span>'
                for x in sl.get("tags", [])
            )
            secs.append(f"""
<section data-transition="slide">
  <h3 style="color:{C['navy']};font-size:1.2em;font-weight:700;text-align:left;margin-bottom:0.5em">
    {esc(sl['title'])}</h3>
  <div style="display:grid;grid-template-columns:1fr 260px;gap:20px">
    <div>{rows_html}</div>
    <div style="background:{C['navy']};border-left:4px solid {C['teal']};padding:10px">
      {right_html}</div>
  </div>
  <div style="margin-top:0.8em">{tag_html}</div>
</section>""")

        elif t == "cards3":
            cards_html = ""
            for card in sl["cards"]:
                clr = C.get(card["color"], C["teal"])
                is_hl = card.get("highlight", False)
                items_html = "".join(
                    f'<li style="font-size:0.72em;margin:5px 0;color:{C["text"]}">▸ {esc(x)}</li>'
                    for x in card["items"]
                )
                note_bg = C["teal"] if is_hl else C["muted"]
                cards_html += f"""
<div class="fragment" style="flex:1;background:#fff;border-top:5px solid {clr};
                              padding:16px;box-shadow:2px 2px 8px rgba(0,0,0,.12)">
  <div style="font-size:2.5em;font-weight:800;color:{clr};text-align:center">{esc(card['num'])}</div>
  <div style="font-size:0.65em;color:{C['muted']};text-align:center;margin-bottom:8px">{esc(card['unit'])}</div>
  <div style="font-weight:700;font-size:0.8em;color:{C['navy']};text-align:center;margin-bottom:10px">
    {esc(card['label'])}</div>
  <ul style="list-style:none;margin:0;padding:0">{items_html}</ul>
  <div style="margin-top:12px;background:{note_bg};color:#fff;
              font-size:0.65em;padding:4px;text-align:center;font-weight:700">
    {esc(card.get('note',''))}</div>
</div>"""
            secs.append(f"""
<section data-background-color="{C['lgray']}" data-transition="slide">
  <h3 style="color:{C['navy']};font-size:1.2em;font-weight:700;text-align:left;margin-bottom:0.3em">
    {esc(sl['title'])}</h3>
  <p style="color:{C['muted']};font-size:0.65em;text-align:left;margin-bottom:0.8em;font-style:italic">
    {esc(sl.get('subtitle',''))}</p>
  <div style="display:flex;gap:16px">{cards_html}</div>
</section>""")

        elif t == "bullets":
            items_html = ""
            for item in sl["items"]:
                is_arrow = str(item).startswith("→")
                clr = C["teal"] if is_arrow else C["text"]
                fw = "700" if is_arrow else "400"
                bullet = "" if is_arrow else "▸&ensp;"
                items_html += f"""
<li class="fragment" style="list-style:none;margin:0.4em 0;
                             color:{clr};font-weight:{fw};font-size:0.82em">
  {bullet}{esc(item)}</li>"""
            tag_html = "".join(
                f'<span style="background:{C["teal"]};color:#fff;padding:2px 10px;'
                f'border-radius:4px;font-size:0.55em;margin-right:6px">{esc(x)}</span>'
                for x in sl.get("tags", [])
            )
            note_html = f'<p style="font-size:0.55em;color:{C["muted"]};font-style:italic;margin-top:0.8em">{esc(sl["note"])}</p>' if sl.get("note") else ""
            secs.append(f"""
<section data-transition="slide">
  <h3 style="color:{C['navy']};font-size:1.2em;font-weight:700;text-align:left;margin-bottom:0.5em">
    {esc(sl['title'])}</h3>
  <ul style="text-align:left;margin:0;padding:0">{items_html}</ul>
  {note_html}
  <div style="margin-top:0.8em">{tag_html}</div>
</section>""")

        elif t == "results":
            bars_html = ""
            for bar in sl["bars"]:
                clr = C.get(bar["color"], C["teal"])
                is_hl = bar.get("highlight", False)
                pct = bar["value"] * 100
                hl_style = f"background:#f0fdfa;border-radius:4px;padding:6px 8px;" if is_hl else "padding:6px 8px;"
                bars_html += f"""
<div class="fragment" style="{hl_style}margin:10px 0">
  <div style="display:flex;align-items:center;gap:12px;margin-bottom:4px">
    <span style="font-size:0.75em;font-weight:{'700' if is_hl else '400'};
                 color:{C['navy']};min-width:130px">{esc(bar['label'])}</span>
    <span style="background:{clr};color:#fff;font-size:0.55em;
                 padding:2px 8px;border-radius:3px">{esc(bar['feat'])}</span>
  </div>
  <div style="display:flex;align-items:center;gap:10px">
    <div style="flex:1;height:22px;background:#e2e8f0;border-radius:3px">
      <div style="width:{pct:.1f}%;height:100%;background:{clr};border-radius:3px"></div>
    </div>
    <span style="font-size:0.8em;font-weight:700;color:{C['navy']};min-width:60px">
      {bar['value']:.4f}</span>
    <span style="background:{clr};color:#fff;font-size:0.55em;
                 padding:2px 8px;border-radius:3px">{esc(bar.get('tag',''))}</span>
  </div>
</div>"""
            secs.append(f"""
<section data-transition="slide">
  <h3 style="color:{C['navy']};font-size:1.2em;font-weight:700;text-align:left;margin-bottom:0.3em">
    {esc(sl['title'])}</h3>
  <p style="color:{C['muted']};font-size:0.65em;text-align:left;margin-bottom:0.8em;font-style:italic">
    {esc(sl.get('subtitle',''))}</p>
  {bars_html}
  <p style="font-size:0.7em;font-weight:700;color:{C['teal']};margin-top:0.8em">
    ▸ 三個特徵集最佳演算法一致：XGBoost + SMOTE</p>
</section>""")

    html = f"""<!DOCTYPE html>
<html lang="zh-TW">
<head>
<meta charset="UTF-8">
<title>生活習慣與肥胖風險預測</title>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.6.0/dist/reveal.css">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.6.0/dist/theme/white.css">
<style>
  :root {{ --r-main-font-size: 30px; }}
  .reveal {{ font-family: "微軟正黑體", "Segoe UI", sans-serif; }}
  .reveal h1,.reveal h2,.reveal h3 {{
    font-family: "微軟正黑體", "Segoe UI", sans-serif;
    text-transform: none; letter-spacing: normal;
  }}
  .reveal section {{ padding: 0 1.2em; }}
  .reveal ul,.reveal ol {{ display:block; margin:0; }}
  .reveal .slides {{ text-align: left; }}
  .reveal .progress {{ color: {C["teal"]}; }}
</style>
</head>
<body>
<div class="reveal"><div class="slides">
{"".join(secs)}
</div></div>
<script src="https://cdn.jsdelivr.net/npm/reveal.js@4.6.0/dist/reveal.js"></script>
<script>
Reveal.initialize({{hash:true,transition:'slide',transitionSpeed:'fast',
  center:false,controls:true,progress:true,slideNumber:'c/t'}});
</script>
</body></html>"""

    with open(out, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"Reveal →  {out}")


# ═══════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════
if __name__ == "__main__":
    gen_pptx(SLIDES,   "outputs/slides.pptx")
    gen_reveal(SLIDES, "outputs/slides_reveal.html")
    print("\nDone.")
